#!/usr/bin/env python3
"""Generate Bedrock Cost Optimization brownbag slides from CI&T template.

Uses the CI&T Slide Presentation Template as the base, clones relevant slides,
modifies content, and adds speaker notes for a 30-minute presentation.
"""

import sys, os, copy
from lxml import etree

# Add venv packages
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', '.venv', 'lib', 'python3.13', 'site-packages'))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

VAULT = os.path.join(os.path.dirname(__file__), '..')
TEMPLATE = os.path.join(VAULT, 'Work', 'Brownbag Sessions', 'CI&T - Slide Presentation Template.pptx')
OUTPUT = os.path.join(VAULT, 'Work', 'Brownbag Sessions', 'Bedrock Cost Optimization', 'Bedrock Cost Optimization Slides.pptx')

# ── Namespace helpers ──
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
}

def qn(tag):
    prefix, local = tag.split(':')
    return f'{{{NSMAP[prefix]}}}{local}'


# ── Slide cloning ──
def clone_slide(prs, src_index):
    """Clone slide at src_index, append to end, return new slide.

    Modifies the spTree in-place so python-pptx's cached reference stays valid.
    """
    src = prs.slides[src_index]
    layout = src.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Modify existing spTree in-place (python-pptx caches spTree reference)
    new_spTree = new_slide.shapes._spTree
    src_spTree = src.shapes._spTree

    # Remove non-structural children from new spTree
    for child in list(new_spTree):
        if not child.tag.endswith('}nvGrpSpPr') and not child.tag.endswith('}grpSpPr'):
            new_spTree.remove(child)

    # Copy shape elements from source
    for child in src_spTree:
        if not child.tag.endswith('}nvGrpSpPr') and not child.tag.endswith('}grpSpPr'):
            new_spTree.append(copy.deepcopy(child))

    # Copy background
    src_bg = src.element.find(qn('p:cSld')).find(qn('p:bg'))
    new_cSld = new_slide.element.find(qn('p:cSld'))
    old_bg = new_cSld.find(qn('p:bg'))
    if src_bg is not None:
        new_bg = copy.deepcopy(src_bg)
        if old_bg is not None:
            new_cSld.remove(old_bg)
        new_cSld.insert(0, new_bg)

    # Copy image relationships
    for rel in src.part.rels.values():
        if 'image' in rel.reltype:
            new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)

    return new_slide


def delete_slide(prs, index):
    """Delete slide at index."""
    rId = prs.slides._sldIdLst[index].get(qn('r:id'))
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[index]


def _get_first_rPr(shape):
    """Extract first run's formatting element from a shape."""
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn('a:rPr'))
            if rPr is not None:
                return copy.deepcopy(rPr)
    return None


def _get_first_pPr(shape):
    """Extract first paragraph's properties element from a shape."""
    tf = shape.text_frame
    for para in tf.paragraphs:
        pPr = para._p.find(qn('a:pPr'))
        if pPr is not None:
            return copy.deepcopy(pPr)
    return None


def set_text(shape, new_text):
    """Replace all text in a shape while preserving formatting."""
    rPr = _get_first_rPr(shape)
    pPr = _get_first_pPr(shape)
    txBody = shape.text_frame._txBody

    # Remove all existing paragraphs
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    # Add single paragraph with text
    new_p = etree.SubElement(txBody, qn('a:p'))
    if pPr is not None:
        new_p.append(pPr)
    new_r = etree.SubElement(new_p, qn('a:r'))
    if rPr is not None:
        new_r.append(copy.deepcopy(rPr))
    t = etree.SubElement(new_r, qn('a:t'))
    t.text = new_text


def set_multiline_text(shape, lines):
    """Set multiple lines in a shape, preserving formatting for all lines."""
    rPr = _get_first_rPr(shape)
    pPr = _get_first_pPr(shape)
    txBody = shape.text_frame._txBody

    # Remove all existing paragraphs
    for p in list(txBody.findall(qn('a:p'))):
        txBody.remove(p)

    # Add new paragraphs
    for line in lines:
        new_p = etree.SubElement(txBody, qn('a:p'))
        if pPr is not None:
            new_p.append(copy.deepcopy(pPr))
        new_r = etree.SubElement(new_p, qn('a:r'))
        if rPr is not None:
            new_r.append(copy.deepcopy(rPr))
        t = etree.SubElement(new_r, qn('a:t'))
        t.text = line


def add_speaker_notes(slide, notes_text):
    """Add or replace speaker notes on a slide."""
    if not slide.has_notes_slide:
        slide.notes_slide  # creates notes slide
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def get_text_shapes(slide):
    """Return list of shapes with text, sorted top-to-bottom."""
    shapes = [s for s in slide.shapes if s.has_text_frame and s.text_frame.text.strip()]
    shapes.sort(key=lambda s: (s.top or 0, s.left or 0))
    return shapes


# ══════════════════════════════════════════════════════════════
#  SLIDE CONTENT DEFINITION
# ══════════════════════════════════════════════════════════════

# Template slide indices for each purpose:
#   8  = Title cover (coral bg + geometric shapes)
#  17  = Agenda 4-topic (maroon bg, big numbers 01-04)
#  19  = Statement/section (navy bg, centered big text)
#  20  = Highlight statement (light blue bg)
#  21  = Numbered section "01." (pink bg)
#  24  = Two-column content (navy bg)
#  27  = Content with title + paragraph (light blue bg)
#  29  = Three-column content (maroon bg)
#  48  = Numbered items x4 (coral bg)
#  81  = Thank you (light blue)

SLIDE_PLAN = [
    # (template_index, content_dict, speaker_notes)
    # 1. Title Cover
    (8, {
        0: "Bedrock Inference\nCost Optimization",
        1: "4 strategies to reduce your AI inference costs on AWS",
        2: "CI&T Engineering",
    }, """Welcome everyone to today's brownbag session. I'm going to walk you through four practical strategies for optimizing inference costs on Amazon Bedrock.

As more of our projects adopt generative AI, inference costs are becoming a significant line item. The good news is there are proven techniques to reduce these costs — often by 50% or more — without sacrificing output quality.

Today's session is structured from simplest to most complex, so you can pick the right starting point for your team's maturity level."""),

    # 2. Agenda
    (17, {
        'topics': [
            ("Prompt Compression", "Pure engineering, no AWS dependency"),
            ("Prompt Caching", "Bedrock feature, config and go"),
            ("Model Routing", "Architecture decision: Haiku vs Sonnet"),
            ("Batch Processing", "Workflow redesign for offline tasks"),
        ]
    }, """Here's our agenda. We'll cover four strategies in order of implementation complexity.

First, Prompt Compression — this is pure engineering that you can do today with no AWS changes. Second, Prompt Caching — a Bedrock feature you just need to enable. Third, Model Routing — this requires an architectural decision about when to use cheaper models. And fourth, Batch Processing — which needs workflow redesign but unlocks significant savings for offline workloads.

After covering all four, we'll look at a comparison table and discuss which strategies combine well together. We'll save about 5 minutes at the end for Q&A."""),

    # 3. Section: The Cost Problem
    (19, {
        0: "The Cost Problem",
        1: "Bedrock charges per token — costs grow linearly with usage",
    }, """Let's start with the problem. Amazon Bedrock uses a pay-per-token pricing model. You pay for both input tokens — what you send to the model — and output tokens — what the model generates.

For Claude models on Bedrock, input tokens are roughly $3 per million for Haiku and $15 per million for Sonnet. Output tokens cost 5x more than input tokens.

As your application scales, and as you add more context through RAG or conversation history, these costs grow linearly. A single complex request with a large system prompt and RAG context can easily consume 10,000+ input tokens. Multiply that by thousands of daily requests and the bill adds up quickly."""),

    # 4. Pricing Model Detail
    (24, {
        0: "How Bedrock Token Pricing Works",
        1: """Key pricing facts:

• You pay for BOTH input and output tokens
• Output tokens cost ~5x more than input
• Claude 3.5 Sonnet: $3/$15 per 1M tokens (in/out)
• Claude 3.5 Haiku: $0.80/$4 per 1M tokens
• Longer prompts = higher cost per request
• RAG context and chat history inflate input tokens

The biggest lever: reduce what you send in.""",
    }, """Let me break down the pricing model. Bedrock charges separately for input and output tokens. Output tokens — what the model generates — are about 5 times more expensive than input tokens.

Looking at the Claude model family on Bedrock: Sonnet costs $3 per million input tokens and $15 per million output tokens. Haiku is much cheaper at $0.80 and $4 respectively.

The key insight here is that the biggest cost driver for most applications is the input side — especially system prompts, RAG context, and conversation history. That's where our optimization strategies will focus.

A typical enterprise RAG application might send 3,000 tokens of system prompt plus 5,000 tokens of retrieved context with every single request. That's 8,000 input tokens before the user even asks their question."""),

    # 5. Strategy 1 Header
    (21, {
        0: "01.",
        1: "Prompt Compression",
        2: "Reduce token count through engineering techniques",
    }, """Our first strategy is Prompt Compression. This is the lowest-hanging fruit because it requires zero AWS configuration changes — it's pure engineering work on your prompt design."""),

    # 6. Strategy 1 Content
    (27, {
        0: "Prompt Compression Techniques",
        1: """5 proven techniques to reduce input tokens:

1. Remove redundant instructions — audit your system prompts for repetition
2. Use concise language — replace verbose phrases with shorter equivalents
3. Compress RAG context — summarize retrieved chunks before injection
4. Structured format — use JSON/YAML instead of verbose natural language
5. Dynamic prompt assembly — only include relevant sections per request

Expected savings: 30-70% reduction in input token cost
Combinability: Stacks with ALL other strategies""",
    }, """Let me walk through five specific techniques for prompt compression.

First, audit your system prompts for redundancy. It's common to find the same instruction repeated in different words across a long prompt. Just deduplicating can save 10-20% of tokens.

Second, use concise language. Instead of "Please make sure to always respond in a professional and courteous manner," you can write "Respond professionally." Same meaning, 80% fewer tokens.

Third, and this is a big one for RAG applications — compress your retrieved context. Instead of injecting full document chunks, use a smaller model to summarize them first. This can cut RAG context by 50-70%.

Fourth, consider structured formats. If you're passing configuration or rules to the model, JSON or YAML is more token-efficient than natural language descriptions.

Fifth, dynamic prompt assembly. Instead of one massive system prompt that covers every scenario, build your prompt dynamically based on the specific request type. Only include the instructions relevant to what the user is asking.

Combined, these techniques typically achieve 30-70% reduction in input token costs, and they stack beautifully with all other strategies we'll discuss."""),

    # 7. Strategy 2 Header — clone of slide 21
    (21, {
        0: "02.",
        1: "Prompt Caching",
        2: "Cache repeated prompt prefixes for ~90% read discount",
    }, """Strategy two is Prompt Caching. This is a Bedrock-native feature that can dramatically reduce costs for applications with repeated prompt prefixes."""),

    # 8. Strategy 2 Content
    (27, {
        0: "How Prompt Caching Works",
        1: """Bedrock caches the prefix of your prompt (system prompt + static context).
Subsequent requests that share the same prefix get a ~90% discount on cached tokens.

How to enable:
• Set cache control headers in your Bedrock API calls
• Mark cache breakpoints in your prompt structure
• Minimum cacheable prefix: 1,024 tokens (Haiku) / 2,048 tokens (Sonnet)

Best for:
• Multi-turn conversations with fixed system prompts
• RAG apps with stable few-shot examples
• Any pattern where the same prefix repeats across requests

Watch out:
• Cache write costs slightly more than regular input
• Cache has a TTL — first request after expiry pays full price
• Must meet minimum token thresholds""",
    }, """Prompt Caching works by storing the prefix of your prompt — typically the system prompt and any static context — on Bedrock's infrastructure. When subsequent requests share the same prefix, you get roughly a 90% discount on those cached tokens.

To enable it, you set cache control headers in your API calls and mark where the cacheable prefix ends. There are minimum thresholds — 1,024 tokens for Haiku and 2,048 for Sonnet.

The ideal use cases are multi-turn conversations where every message includes the same system prompt, and RAG applications that include fixed few-shot examples in every request.

One gotcha: the initial cache write costs slightly more than a regular input — about 25% more for Claude models. So you want to make sure the prefix actually gets reused enough to pay back that investment. The cache also has a TTL, typically 5 minutes, so it works best for patterns with frequent requests.

The math works out well for most production applications. If you have a 3,000-token system prompt and handle 100 requests per minute, you're paying full price once every 5 minutes and getting 90% off for the other 499 requests in that window."""),

    # 9. Strategy 3 Header
    (21, {
        0: "03.",
        1: "Model Routing",
        2: "Route simple tasks to cheaper models automatically",
    }, """Strategy three is Model Routing — using the right model for the right task. This is where we start making architectural decisions."""),

    # 10. Strategy 3 Content
    (27, {
        0: "Intelligent Model Routing",
        1: """Core idea: Not every request needs your most powerful (expensive) model.
Route simple tasks to Haiku, complex ones to Sonnet.

Two approaches:

1. Bedrock Intelligent Prompt Routing (IPR)
   • AWS-managed routing between model pairs
   • Automatic complexity classification
   • No code changes needed

2. Custom routing logic
   • Build a classifier (rule-based or ML) to categorize requests
   • Map categories to models: simple → Haiku, complex → Sonnet
   • More control, but more engineering effort

Savings: 60-70% when simple tasks dominate your traffic
Tension with compression: over-compressed prompts + small models = quality risk""",
    }, """The core idea behind model routing is simple: not every request needs your most expensive model. If 70% of your incoming requests are straightforward — like classification, simple Q&A, or summarization — you can route those to Haiku at a fraction of the cost, and reserve Sonnet for the complex reasoning tasks.

There are two main approaches. First, Bedrock's built-in Intelligent Prompt Routing, or IPR. This is an AWS-managed feature that automatically classifies request complexity and routes between model pairs. It's the easiest to implement — essentially zero code changes.

Second, you can build custom routing logic. This gives you more control — you define the rules for what counts as "simple" versus "complex" based on your domain knowledge. You might use keyword matching, request metadata, or even a lightweight ML classifier.

The savings can be dramatic — 60-70% cost reduction when simple tasks make up the majority of your traffic, which is the case for most enterprise applications.

One important tension to be aware of: model routing can conflict with aggressive prompt compression. If you over-compress your prompts AND route to a smaller model, the model might not have enough context to produce good results. You need to find the right balance."""),

    # 11. Strategy 4 Header
    (21, {
        0: "04.",
        1: "Batch Processing",
        2: "Queue offline tasks for batch pricing discounts",
    }, """Our fourth and final strategy is Batch Processing — redesigning your workflow to take advantage of batch pricing for tasks that don't need real-time responses."""),

    # 12. Strategy 4 Content
    (27, {
        0: "Batch Processing for Offline Workloads",
        1: """Bedrock Batch Inference lets you submit large batches of requests
and get results asynchronously at a discounted price.

Ideal workloads:
• Document classification and tagging
• Bulk summarization
• Test data generation
• Evaluation and benchmarking
• Content moderation at scale

How it works:
1. Prepare a JSONL file with all your requests
2. Submit to Bedrock Batch API
3. Results delivered asynchronously (typically within hours)
4. Pay batch pricing (significant discount vs real-time)

Not suitable for: real-time chat, interactive applications, low-latency requirements

Combines well with: Prompt Compression (compress, then batch)""",
    }, """Batch Processing is about identifying workloads that don't need real-time responses and batching them together for cost savings.

Bedrock's Batch Inference API lets you submit thousands of requests in a single JSONL file and get results asynchronously. The pricing is significantly discounted compared to real-time inference.

The best candidates for batching are tasks like document classification, bulk summarization, test data generation, evaluation runs, and content moderation. Basically anything where the user isn't sitting there waiting for a response.

The workflow is straightforward: prepare your requests in JSONL format, submit to the batch API, and collect results when they're ready — typically within a few hours depending on volume.

The key constraint is that this doesn't work for real-time or interactive use cases. You need to redesign your workflow to be asynchronous.

Batch processing combines well with prompt compression — you compress your prompts first, then submit the compressed versions in batch. However, the combination with prompt caching doesn't make sense since batch requests don't benefit from the real-time cache."""),

    # 13. Comparison Table
    (48, {
        'items': [
            ("01", "Prompt Compression", "30-70% savings\n⭐ Low complexity\nStacks with everything"),
            ("02", "Prompt Caching", "~90% cache discount\n⭐ Low complexity\nBest after compression"),
            ("03", "Model Routing", "60-70% savings\n⭐⭐ Medium complexity\nWatch compression tension"),
            ("04", "Batch Processing", "Batch discount\n⭐⭐ Medium complexity\nOffline tasks only"),
        ]
    }, """Let's put all four strategies side by side.

Prompt Compression offers 30 to 70% input cost savings with low implementation complexity. It's the most universally applicable and stacks with every other strategy.

Prompt Caching gives you roughly 90% discount on cached token reads, also with low complexity. It works best when applied AFTER compression — you compress your prompts first, then cache the compressed versions for a double benefit.

Model Routing can save 60 to 70% when simple tasks dominate your traffic. It requires medium complexity — you need to make architectural decisions about classification and model selection. Be careful combining it with aggressive compression.

Batch Processing gives you batch pricing discounts for offline workloads. Medium complexity because you need to redesign workflows to be asynchronous. Only applicable for non-real-time use cases.

My recommendation: start with Prompt Compression — it's the easiest win and benefits everything else. Then add Caching. Consider Routing and Batching based on your specific workload patterns."""),

    # 14. Combination Strategies
    (20, {
        0: "Strategy Combinations\nThat Work",
        1: """✅ Compression + Caching: compress first, then cache — double savings
✅ Compression + Batching: compress prompts, submit in batch
⚠️ Compression + Routing: tension — over-compressed + small model = quality risk

Start simple → layer on complexity as you measure impact""",
    }, """Before we wrap up, let me highlight how these strategies combine.

The golden combination is Compression plus Caching. You compress your prompts first to reduce token count, then cache the compressed version. This gives you both the upfront token reduction AND the 90% cache read discount on the remaining tokens. It's multiplicative savings.

Compression plus Batching also works well for offline workloads. Compress your prompts, then submit the smaller versions in batch.

The combination to be cautious about is Compression plus Routing. If you aggressively compress your prompts and then route to a smaller, cheaper model like Haiku, you might strip away too much context for the smaller model to handle correctly. This can lead to quality degradation and retry loops that actually increase costs.

My recommended adoption path: Start with Prompt Compression since it requires no infrastructure changes and benefits everything else. Once that's in place, enable Prompt Caching for your highest-volume endpoints. Then evaluate Model Routing based on your traffic patterns. Finally, identify batch-eligible workloads and migrate them to the Batch API.

Measure the impact at each step before adding the next layer of complexity."""),

    # 15. Thank You + Q&A
    (81, {
        0: "THANK YOU.",
        1: "ciandt.com",
    }, """That's the end of the prepared content. Let me do a quick recap:

We covered four strategies for Bedrock cost optimization, ordered from simple to complex:
1. Prompt Compression — engineering your prompts to use fewer tokens
2. Prompt Caching — caching repeated prefixes for massive read discounts
3. Model Routing — using cheaper models for simple tasks
4. Batch Processing — batching offline workloads for discounted pricing

The key takeaway: start with compression, layer on caching, then evaluate routing and batching based on your specific workload.

I'm happy to take questions now. If you want to dive deeper into any of these strategies for your specific project, feel free to reach out after the session as well."""),
]


# ══════════════════════════════════════════════════════════════
#  MAIN BUILD
# ══════════════════════════════════════════════════════════════

def build_presentation():
    prs = Presentation(TEMPLATE)
    original_count = len(prs.slides)
    print(f"Template loaded: {original_count} slides")

    new_slides = []

    for template_idx, content, notes in SLIDE_PLAN:
        new_slide = clone_slide(prs, template_idx)
        new_slides.append((new_slide, content, notes))

    # Now populate content in new slides
    for slide, content, notes in new_slides:
        shapes = get_text_shapes(slide)

        if 'topics' in content:
            # Agenda slide — find topic+description shapes
            # Template agenda: "01" "02" "03" "04" then "Topic 1\nBrief desc" x4
            topic_shapes = [s for s in shapes if s.text_frame.text.startswith('Topic')]
            for i, ts in enumerate(topic_shapes):
                if i < len(content['topics']):
                    name, desc = content['topics'][i]
                    set_multiline_text(ts, [name, desc])

        elif 'items' in content:
            # Numbered items slide (comparison table)
            # Has "01" "02" "03" "04" + title shapes + description shapes
            title_shapes = [s for s in shapes
                           if s.text_frame.text.startswith('Title for this number')
                           or s.text_frame.text.startswith('Description of this number')]
            num_shapes = [s for s in shapes if s.text_frame.text.strip() in ('01', '02', '03', '04')]

            # Group shapes by vertical position to match items
            # Each item has: number, title+desc
            all_item_shapes = [s for s in shapes
                              if s.text_frame.text.strip() not in ('Insert here your title. Use up to 2 lines of text.',)]

            # Simpler approach: find shapes by content pattern
            for s in shapes:
                txt = s.text_frame.text.strip()
                if txt == 'Insert here your title. Use up to 2 lines of text.':
                    set_text(s, "Strategy Comparison")
                    continue
                for i, (num, title, desc) in enumerate(content['items']):
                    expected_num = f"0{i+1}"
                    if txt == expected_num:
                        set_text(s, num)
                        break
                    if txt.startswith('Title for this number') or txt.startswith('Description of this number'):
                        # Match by position
                        pass

            # Fallback: modify the title/desc shapes in order
            td_shapes = [s for s in shapes
                        if 'Title for this number' in s.text_frame.text
                        or 'Description of this number' in s.text_frame.text]
            td_shapes.sort(key=lambda s: (s.left or 0))
            for i, (num, title, desc) in enumerate(content['items']):
                if i < len(td_shapes):
                    set_multiline_text(td_shapes[i], [title, "", desc])

            # Set the main title
            for s in shapes:
                if 'Insert here your title' in s.text_frame.text:
                    set_text(s, "Strategy Comparison")
                    break

        else:
            # Generic content: map by index in content dict
            for idx, new_text in content.items():
                if isinstance(idx, int) and idx < len(shapes):
                    if '\n' in new_text:
                        set_multiline_text(shapes[idx], new_text.split('\n'))
                    else:
                        set_text(shapes[idx], new_text)

        # Add speaker notes
        add_speaker_notes(slide, notes)

    # Delete all original template slides (they're at the beginning)
    print(f"Deleting {original_count} original template slides...")
    for _ in range(original_count):
        delete_slide(prs, 0)

    print(f"Final presentation: {len(prs.slides)} slides")

    # Save
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved to: {OUTPUT}")


if __name__ == '__main__':
    build_presentation()
